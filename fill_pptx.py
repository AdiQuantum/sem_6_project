"""
Script to complete the BTechPresentation-MidsemTemplate with MERN E-Commerce project content.
"""
from pptx import Presentation
from pptx.util import Pt
from lxml import etree
import copy

def xml_escape(text):
    """Escape special XML characters in text."""
    return (text
        .replace('&', '&amp;')
        .replace('<', '&lt;')
        .replace('>', '&gt;')
        .replace('"', '&quot;')
        .replace("'", '&apos;')
    )

# ── Student info ──────────────────────────────────────────────────────────────
STUDENT_NAME       = "Aditya"          # ← change if needed
ROLL_NO            = "Your Roll No"    # ← change to your roll number
SUPERVISOR         = "Your Supervisor" # ← change to your supervisor's name
PROJECT_TITLE      = "Online Shopping Portal using MERN Stack"
# ─────────────────────────────────────────────────────────────────────────────

NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

def make_heading_para(text):
    """Return XML <a:p> element for a bold heading paragraph."""
    t = xml_escape(text)
    p_xml = f'''<a:p xmlns:a="{NS}">
  <a:r>
    <a:rPr lang="en-IN" b="1" dirty="0">
      <a:latin typeface="Calibri (Body)"/>
    </a:rPr>
    <a:t>{t}</a:t>
  </a:r>
</a:p>'''
    return etree.fromstring(p_xml)

def make_bullet_para(text, level=1):
    """Return XML <a:p> element for a bullet-point paragraph."""
    t = xml_escape(text)
    p_xml = f'''<a:p xmlns:a="{NS}">
  <a:pPr lvl="{level}">
    <a:buFont typeface="Wingdings" panose="05000000000000000000" pitchFamily="2" charset="2"/>
    <a:buChar char="&#167;"/>
  </a:pPr>
  <a:r>
    <a:rPr lang="en-IN" dirty="0">
      <a:latin typeface="Calibri (Body)"/>
    </a:rPr>
    <a:t>{t}</a:t>
  </a:r>
</a:p>'''
    return etree.fromstring(p_xml)

def make_empty_para(indent=False):
    """Return an empty paragraph."""
    if indent:
        p_xml = f'''<a:p xmlns:a="{NS}">
  <a:pPr marL="0" indent="0"><a:buNone/></a:pPr>
  <a:endParaRPr lang="en-IN" dirty="0">
    <a:latin typeface="Calibri (Body)"/>
  </a:endParaRPr>
</a:p>'''
    else:
        p_xml = f'''<a:p xmlns:a="{NS}">
  <a:endParaRPr lang="en-IN" dirty="0">
    <a:latin typeface="Calibri (Body)"/>
  </a:endParaRPr>
</a:p>'''
    return etree.fromstring(p_xml)

def set_text_box(shape, text):
    """Replace the text inside a simple text-box shape."""
    tf = shape.text_frame
    tf.paragraphs[0].runs[0].text = text

def replace_txbody(shape, paragraphs_xml):
    """Replace all content paragraphs in a placeholder with the given list of XML elements."""
    txBody = shape.text_frame._txBody
    # Remove existing <a:p> children
    for p in txBody.findall(f'{{{NS}}}p'):
        txBody.remove(p)
    # Append new paragraphs
    for p in paragraphs_xml:
        txBody.append(p)

# ── Load template ─────────────────────────────────────────────────────────────
prs = Presentation("BTechPresentation-MidsemTemplate (2).pptx")

# =============================================================================
# SLIDE 1  –  Introduction & Motivation
# =============================================================================
slide1 = prs.slides[0]

# 1a. Title text box
for shape in slide1.shapes:
    if shape.name == "Title 1":
        tf = shape.text_frame
        # Clear and set
        for para in tf.paragraphs:
            for run in para.runs:
                run.text = ""
        # Set using raw XML approach
        txBody = tf._txBody
        for p in txBody.findall(f'{{{NS}}}p'):
            txBody.remove(p)
        title_p = etree.fromstring(f'''<a:p xmlns:a="{NS}">
  <a:r>
    <a:rPr lang="en-IN" b="1" dirty="0">
      <a:latin typeface="Calibri (Body)"/>
    </a:rPr>
    <a:t>Project Title: {PROJECT_TITLE}</a:t>
  </a:r>
</a:p>''')
        txBody.append(title_p)
        break

# 1b. Content placeholder
intro_bullets = [
    "An E-Commerce web application built using the MERN Stack (MongoDB, Express.js, React.js, Node.js)",
    "Enables users to browse products, add to cart, place orders, and make secure online payments",
    "Includes an Admin Panel for managing products, orders, and user accounts",
]
motivation_bullets = [
    "India's e-commerce market is growing rapidly — online shopping is becoming the norm for consumers",
    "Existing platforms like Amazon and Flipkart are complex and proprietary; building our own demonstrates full-stack competency",
    "The MERN stack is one of the most in-demand technology stacks in the industry, making this project highly relevant for placements",
]

paras = []
paras.append(make_heading_para("Introduction:"))
for b in intro_bullets:
    paras.append(make_bullet_para(b))
paras.append(make_empty_para())
paras.append(make_heading_para("Motivation:"))
for b in motivation_bullets:
    paras.append(make_bullet_para(b))
paras.append(make_empty_para(indent=True))

for shape in slide1.shapes:
    if shape.name == "Content Placeholder 2":
        replace_txbody(shape, paras)
        break

# 1c. Footer group – Roll No, Name, Supervisor
for shape in slide1.shapes:
    if shape.shape_type == 6 and shape.name == "Group 1":
        for child in shape.shapes:
            if child.has_text_frame:
                text = child.text_frame.paragraphs[0].text.strip()
                tf = child.text_frame
                if text.startswith("Roll No"):
                    txBody = tf._txBody
                    for p in txBody.findall(f'{{{NS}}}p'):
                        txBody.remove(p)
                    txBody.append(etree.fromstring(f'<a:p xmlns:a="{NS}"><a:r><a:rPr lang="en-IN" b="1" dirty="0"><a:latin typeface="Calibri (Body)"/></a:rPr><a:t>Roll No: {ROLL_NO}</a:t></a:r></a:p>'))
                elif text.startswith("Name"):
                    txBody = tf._txBody
                    for p in txBody.findall(f'{{{NS}}}p'):
                        txBody.remove(p)
                    txBody.append(etree.fromstring(f'<a:p xmlns:a="{NS}"><a:r><a:rPr lang="en-IN" b="1" dirty="0"><a:latin typeface="Calibri (Body)"/></a:rPr><a:t>Name: {STUDENT_NAME}</a:t></a:r></a:p>'))
                elif text.startswith("Supervisor"):
                    txBody = tf._txBody
                    for p in txBody.findall(f'{{{NS}}}p'):
                        txBody.remove(p)
                    txBody.append(etree.fromstring(f'<a:p xmlns:a="{NS}"><a:r><a:rPr lang="en-IN" b="1" dirty="0"><a:latin typeface="Calibri (Body)"/></a:rPr><a:t>Supervisor: {SUPERVISOR}</a:t></a:r></a:p>'))
        break

# =============================================================================
# SLIDE 2  –  Existing Work / Literature Review & Gaps
# =============================================================================
slide2 = prs.slides[1]

existing_work_bullets = [
    "Amazon and Flipkart are feature-rich platforms but are not open-source and cannot be studied or customised easily",
    "OpenCart and WooCommerce are open-source e-commerce frameworks but are built on PHP, not JavaScript; learning curve differs",
    "Academic MERN e-commerce projects exist on GitHub but often lack payment integration, cloud deployment, and admin functionality",
]
gaps_bullets = [
    "Most open-source projects lack integrated payment gateways (Razorpay / Stripe) for the Indian market",
    "Existing academic projects are rarely deployed to production-grade cloud environments (Vercel + MongoDB Atlas)",
    "Few projects combine both a functional user-facing storefront AND a complete Admin Panel in a single MERN application",
]

paras2 = []
paras2.append(make_heading_para("Existing Work / Literature Review:"))
for b in existing_work_bullets:
    paras2.append(make_bullet_para(b))
paras2.append(make_empty_para())
paras2.append(make_heading_para("Gaps in the Existing Solutions / Work:"))
for b in gaps_bullets:
    paras2.append(make_bullet_para(b))
paras2.append(make_empty_para(indent=True))

for shape in slide2.shapes:
    if shape.name == "Content Placeholder 2":
        replace_txbody(shape, paras2)
        break

# Footer for slide 2
for shape in slide2.shapes:
    if shape.shape_type == 6 and shape.name == "Group 1":
        for child in shape.shapes:
            if child.has_text_frame:
                text = child.text_frame.paragraphs[0].text.strip()
                tf = child.text_frame
                if text.startswith("Roll No"):
                    txBody = tf._txBody
                    for p in txBody.findall(f'{{{NS}}}p'):
                        txBody.remove(p)
                    txBody.append(etree.fromstring(f'<a:p xmlns:a="{NS}"><a:r><a:rPr lang="en-IN" b="1" dirty="0"><a:latin typeface="Calibri (Body)"/></a:rPr><a:t>Roll No: {ROLL_NO}</a:t></a:r></a:p>'))
                elif text.startswith("Name"):
                    txBody = tf._txBody
                    for p in txBody.findall(f'{{{NS}}}p'):
                        txBody.remove(p)
                    txBody.append(etree.fromstring(f'<a:p xmlns:a="{NS}"><a:r><a:rPr lang="en-IN" b="1" dirty="0"><a:latin typeface="Calibri (Body)"/></a:rPr><a:t>Name: {STUDENT_NAME}</a:t></a:r></a:p>'))
                elif text.startswith("Project Title"):
                    txBody = tf._txBody
                    for p in txBody.findall(f'{{{NS}}}p'):
                        txBody.remove(p)
                    txBody.append(etree.fromstring(f'<a:p xmlns:a="{NS}"><a:r><a:rPr lang="en-IN" b="1" dirty="0"><a:latin typeface="Calibri (Body)"/></a:rPr><a:t>Project Title: {PROJECT_TITLE}</a:t></a:r></a:p>'))
        break

# =============================================================================
# SLIDE 3  –  Details of Proposed Work / Results Achieved
# =============================================================================
slide3 = prs.slides[2]

proposed_bullets = [
    "Full-stack MERN e-commerce application with React.js frontend, Node.js/Express.js backend, and MongoDB Atlas database",
    "User-facing features: Registration & Login (JWT auth), Product Catalogue with search & filter, Shopping Cart, Order Placement, Payment (Razorpay/Stripe)",
    "Admin Panel: Add/Edit/Delete products (with Cloudinary image upload), Manage Orders & update status, View registered users",
    "Deployed on Vercel (frontend + serverless backend) with MongoDB Atlas as the cloud database",
    "Responsive UI built with React.js + Tailwind CSS; REST API communication between frontend and backend",
]
results_bullets = [
    "Functional web application successfully deployed at a Vercel public URL",
    "End-to-end order flow working: user can register → browse → add to cart → checkout → pay → view order status",
    "Admin dashboard allows real-time product and order management with cloud image storage via Cloudinary",
]

paras3 = []
paras3.append(make_heading_para("Details of Proposed Work / Solutions:"))
for b in proposed_bullets:
    paras3.append(make_bullet_para(b))
paras3.append(make_empty_para())
paras3.append(make_heading_para("Result Achieved (if any):"))
for b in results_bullets:
    paras3.append(make_bullet_para(b))
paras3.append(make_empty_para(indent=True))

for shape in slide3.shapes:
    if shape.name == "Content Placeholder 2":
        replace_txbody(shape, paras3)
        break

# Footer for slide 3
for shape in slide3.shapes:
    if shape.shape_type == 6 and shape.name == "Group 1":
        for child in shape.shapes:
            if child.has_text_frame:
                text = child.text_frame.paragraphs[0].text.strip()
                tf = child.text_frame
                if text.startswith("Roll No"):
                    txBody = tf._txBody
                    for p in txBody.findall(f'{{{NS}}}p'):
                        txBody.remove(p)
                    txBody.append(etree.fromstring(f'<a:p xmlns:a="{NS}"><a:r><a:rPr lang="en-IN" b="1" dirty="0"><a:latin typeface="Calibri (Body)"/></a:rPr><a:t>Roll No: {ROLL_NO}</a:t></a:r></a:p>'))
                elif text.startswith("Name"):
                    txBody = tf._txBody
                    for p in txBody.findall(f'{{{NS}}}p'):
                        txBody.remove(p)
                    txBody.append(etree.fromstring(f'<a:p xmlns:a="{NS}"><a:r><a:rPr lang="en-IN" b="1" dirty="0"><a:latin typeface="Calibri (Body)"/></a:rPr><a:t>Name: {STUDENT_NAME}</a:t></a:r></a:p>'))
                elif text.startswith("Project Title"):
                    txBody = tf._txBody
                    for p in txBody.findall(f'{{{NS}}}p'):
                        txBody.remove(p)
                    txBody.append(etree.fromstring(f'<a:p xmlns:a="{NS}"><a:r><a:rPr lang="en-IN" b="1" dirty="0"><a:latin typeface="Calibri (Body)"/></a:rPr><a:t>Project Title: {PROJECT_TITLE}</a:t></a:r></a:p>'))
        break

# =============================================================================
# Save
# =============================================================================
out_path = "BTechPresentation-Completed.pptx"
prs.save(out_path)
print(f"DONE! Saved completed presentation as: {out_path}")
